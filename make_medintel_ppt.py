"""
MedIntel — Project Overview deck generator (Software Architecture, Sem 5).
Run:  python make_medintel_ppt.py
Out:  MedIntel_Project_Overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# ---------------------------------------------------------------- design tokens
INK        = RGBColor(0x0F, 0x17, 0x2A)
SLATE      = RGBColor(0x47, 0x55, 0x69)
MUTED      = RGBColor(0x94, 0xA3, 0xB8)
TEAL       = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK  = RGBColor(0x0F, 0x76, 0x6E)
TEAL_DEEP  = RGBColor(0x13, 0x4E, 0x4A)
MINT       = RGBColor(0xCC, 0xFB, 0xF1)
MINT_SOFT  = RGBColor(0xEC, 0xFD, 0xF5)
LIGHT      = RGBColor(0xF1, 0xF5, 0xF9)
LINE       = RGBColor(0xE2, 0xE8, 0xF0)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
AMBER      = RGBColor(0xB4, 0x53, 0x09)
AMBER_BG   = RGBColor(0xFF, 0xF7, 0xED)
ROSE       = RGBColor(0xBE, 0x12, 0x3C)
INDIGO     = RGBColor(0x43, 0x38, 0xCA)

FONT = "Segoe UI"
W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width, prs.slide_height = In(W), In(H)
BLANK = prs.slide_layouts[6]

_num = {"n": 1}  # slide 1 is the unnumbered title slide; footers then match real slide numbers


# ---------------------------------------------------------------- primitives
def txbox(slide, l, t, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(In(l), In(t), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size=14, color=INK, bold=False, first=False,
         space_before=0, space_after=6, align=PP_ALIGN.LEFT, line=1.25, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.italic = Pt(size), bold, italic
    r.font.color.rgb, r.font.name = color, FONT
    return p


def _hanging_bullet(p, char, color, pct, mar):
    """Real DrawingML bullet + hanging indent, so wrapped lines align under the text."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(In(mar))))
    pPr.set("indent", str(int(-In(mar))))
    for xml in (
        '<a:buClr %s><a:srgbClr val="%s"/></a:buClr>' % (nsdecls("a"), str(color)),
        '<a:buSzPct %s val="%d"/>' % (nsdecls("a"), int(pct * 1000)),
        '<a:buFont %s typeface="Arial" pitchFamily="34" charset="0"/>' % nsdecls("a"),
        '<a:buChar %s char="%s"/>' % (nsdecls("a"), char),
    ):
        pPr.append(parse_xml(xml))


def bullet(tf, text, size=14, color=INK, mark="▪", mark_color=TEAL,
           first=False, space_after=9, indent=0, bold=False, mar=0.2):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = 1.22
    r = p.add_run()
    r.text = text
    r.font.size, r.font.name, r.font.bold = Pt(size), FONT, bold
    r.font.color.rgb = color
    if mark:
        _hanging_bullet(p, mark, mark_color, 80 if mark == "▪" else 95, mar)
    return p


def rich(tf, parts, size=14, first=False, space_after=9, line=1.22):
    """parts = [(text, color, bold), ...] rendered in one paragraph."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_after = Pt(space_after)
    p.line_spacing = line
    for text, color, bold in parts:
        r = p.add_run()
        r.text = text
        r.font.size, r.font.name = Pt(size), FONT
        r.font.bold, r.font.color.rgb = bold, color
    return p


def box(slide, l, t, w, h, fill=WHITE, edge=LINE, radius=0.08, edge_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape, In(l), In(t), In(w), In(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if edge is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = edge
        s.line.width = Pt(edge_w)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    return s


def label(shape, text, size=13, color=INK, bold=True, align=PP_ALIGN.CENTER, sub=None,
          sub_size=10, sub_color=SLATE):
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = In(0.08)
    tf.margin_top = tf.margin_bottom = In(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.1
    r = p.add_run()
    r.text = text
    r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
    r.font.color.rgb = color
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = align
        p2.space_before = Pt(2)
        p2.line_spacing = 1.05
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size, r2.font.name = Pt(sub_size), FONT
        r2.font.color.rgb = sub_color
    return shape


def arrow(slide, l, t, w, h, direction="right", color=MINT, edge=None):
    shp = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
           "left": MSO_SHAPE.LEFT_ARROW, "up": MSO_SHAPE.UP_ARROW}[direction]
    s = slide.shapes.add_shape(shp, In(l), In(t), In(w), In(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if edge is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = edge
    s.shadow.inherit = False
    return s


# ---------------------------------------------------------------- slide chrome
def slide(title, kicker=None, notes=None, dark=False):
    s = prs.slides.add_slide(BLANK)
    _num["n"] += 1

    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = TEAL_DEEP if dark else WHITE

    accent = box(s, 0.62, 0.46, 0.055, 0.42, fill=TEAL, edge=None, radius=0.45)
    accent.text_frame.text = ""

    top = 0.42
    if kicker:
        tf = txbox(s, 0.82, 0.40, 10.5, 0.24)
        para(tf, kicker.upper(), size=9.5, color=(MINT if dark else TEAL_DARK),
             bold=True, first=True, space_after=0)
        top = 0.62
    tf = txbox(s, 0.82, top, 11.6, 0.62)
    para(tf, title, size=27, color=(WHITE if dark else INK), bold=True,
         first=True, space_after=0, line=1.05)

    # footer
    fl = box(s, 0.62, 6.92, 12.1, 0.012,
             fill=(RGBColor(0x2A, 0x6C, 0x67) if dark else LINE), edge=None, shape=MSO_SHAPE.RECTANGLE)
    fl.text_frame.text = ""
    ft = txbox(s, 0.62, 7.02, 6.0, 0.24)
    para(ft, "MedIntel  ·  AI-Powered Medical Assistant", size=9,
         color=(MINT if dark else MUTED), first=True, space_after=0)
    fn = txbox(s, 10.7, 7.02, 2.02, 0.24)
    para(fn, f"{_num['n']:02d}", size=9, color=(MINT if dark else MUTED),
         first=True, space_after=0, align=PP_ALIGN.RIGHT)

    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def table(slide_, headers, rows, l, t, w, col_w, fs=11.5, hfs=11,
          row_h=0.42, head_h=0.38, zebra=True, bold_col0=True):
    n_r, n_c = len(rows) + 1, len(headers)
    tbl = slide_.shapes.add_table(n_r, n_c, In(l), In(t), In(w), In(head_h + row_h * len(rows))).table
    tbl.first_row = False
    tbl.horz_banding = False
    total = sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(int(In(w) * cw / total))
    tbl.rows[0].height = In(head_h)
    for r in range(1, n_r):
        tbl.rows[r].height = In(row_h)

    def style(cell, text, size, color, bold, fill, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.margin_left = cell.margin_right = In(0.11)
        cell.margin_top = cell.margin_bottom = In(0.05)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = text
        r.font.size, r.font.bold, r.font.name = Pt(size), bold, FONT
        r.font.color.rgb = color

    for c, htext in enumerate(headers):
        style(tbl.cell(0, c), htext, hfs, WHITE, True, TEAL_DARK)
    for r, row in enumerate(rows, start=1):
        fill = WHITE if (r % 2 or not zebra) else LIGHT
        for c, val in enumerate(row):
            is0 = (c == 0 and bold_col0)
            style(tbl.cell(r, c), val, fs, INK if is0 else SLATE, is0, fill)
    return tbl


# ================================================================ 01 TITLE
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = TEAL_DEEP
for i, (x, y, d, col) in enumerate([(10.4, -1.5, 5.2, RGBColor(0x11, 0x5E, 0x59)),
                                    (11.6, 4.2, 4.0, RGBColor(0x0E, 0x57, 0x53)),
                                    (-1.2, 5.0, 3.4, RGBColor(0x11, 0x5E, 0x59))]):
    c = box(s, x, y, d, d, fill=col, edge=None, shape=MSO_SHAPE.OVAL)
    c.text_frame.text = ""

pill = box(s, 1.0, 1.62, 2.62, 0.36, fill=RGBColor(0x11, 0x5E, 0x59), edge=TEAL, radius=0.5)
label(pill, "SOFTWARE ARCHITECTURE  ·  SEM 5", size=9.5, color=MINT, bold=True)

tf = txbox(s, 1.0, 2.22, 10.4, 1.4)
para(tf, "MedIntel", size=62, color=WHITE, bold=True, first=True, space_after=2, line=1.0)
tf = txbox(s, 1.0, 3.38, 10.2, 0.9)
para(tf, "An AI-Powered Medical Assistant for Symptom Triage,\nHealth Records and Care Adherence",
     size=19, color=MINT, first=True, space_after=0, line=1.28)

rule = box(s, 1.0, 4.44, 1.5, 0.035, fill=TEAL, edge=None, shape=MSO_SHAPE.RECTANGLE)
rule.text_frame.text = ""

tf = txbox(s, 1.0, 4.78, 10.4, 0.5)
para(tf, "Layered Architecture  ·  React.js  ·  Node.js + Express  ·  MongoDB  ·  Gemini / OpenAI  ·  REST",
     size=12.5, color=RGBColor(0x99, 0xE6, 0xDC), first=True, space_after=0)

warn = box(s, 1.0, 5.5, 6.9, 0.52, fill=RGBColor(0x11, 0x5E, 0x59), edge=RGBColor(0x2A, 0x6C, 0x67), radius=0.12)
label(warn, "⚠   Clinical decision SUPPORT — not a diagnosis, not a doctor.",
      size=11.5, color=MINT, bold=False, align=PP_ALIGN.CENTER)

tf = txbox(s, 1.0, 6.42, 10.4, 0.4)
para(tf, "Presented by  •  Member 1  •  Member 2  •  Member 3  •  Member 4",
     size=11, color=RGBColor(0x7F, 0xD1, 0xC6), first=True, space_after=0)
s.notes_slide.notes_text_frame.text = (
    "MEMBER 1 opens.\n\n'Good morning. Our project is MedIntel, an AI-powered medical assistant.\n"
    "The one line to remember: MedIntel is clinical decision SUPPORT, not a diagnosis and not a "
    "replacement for a doctor. That constraint is not a disclaimer we bolted on at the end — you "
    "will see it drive our actual architecture.'")

# ================================================================ 02 PROBLEM
s = slide("The Problem We Are Solving", "Motivation",
          notes="MEMBER 1.\nWalk the four gaps. Land the last line: these are four separate problems today, "
                "and no single tool joins them up. Do not invent statistics if asked — say it is based on "
                "common observation and secondary reading.")
cards = [
    ("\U0001F50D", "Unreliable self-diagnosis", "People search symptoms online and get contradictory, "
     "fear-inducing results with no personal context.", TEAL),
    ("⏱", "Pressure on primary care", "Minor, self-limiting complaints occupy consultation slots, "
     "lengthening queues for those who truly need them.", INDIGO),
    ("\U0001F4C1", "Scattered health records", "Prescriptions on paper, labs as PDFs in e-mail. "
     "History is lost exactly when a doctor needs it.", AMBER),
    ("\U0001F48A", "Poor medication adherence", "Missed and mistimed doses silently undermine "
     "otherwise correct treatment.", ROSE),
]
for i, (icon, head, body, col) in enumerate(cards):
    x = 0.62 + (i % 2) * 6.15
    y = 1.68 + (i // 2) * 2.24
    c = box(s, x, y, 5.85, 1.9, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    bar = box(s, x, y, 0.05, 1.9, fill=col, edge=None, radius=0.4)
    bar.text_frame.text = ""
    tf = txbox(s, x + 0.32, y + 0.32, 5.3, 1.4)
    para(tf, f"{icon}  {head}", size=15, color=INK, bold=True, first=True, space_after=8)
    para(tf, body, size=12.5, color=SLATE, line=1.32, space_after=0)

band = box(s, 0.62, 6.18, 12.1, 0.6, fill=MINT_SOFT, edge=MINT, radius=0.1)
label(band, "These are four disconnected problems today — MedIntel joins them into one guided, "
            "record-keeping loop.", size=12.5, color=TEAL_DARK, bold=False)

# ================================================================ 03 SCOPE
s = slide("What MedIntel Is — and What It Is Not", "Scope & Safety Boundary",
          notes="MEMBER 1.\nThis slide wins marks. Say: 'Defining what the system refuses to do is itself an "
                "architectural constraint. Because we can never guarantee a language model is correct, we "
                "never place it alone on a safety-critical path — Member 3 will show the deterministic "
                "rules engine that sits in front of it.'")
head_is = box(s, 0.62, 1.55, 5.85, 0.46, fill=MINT_SOFT, edge=MINT, radius=0.14)
label(head_is, "✓   IT IS", size=13, color=TEAL_DARK, align=PP_ALIGN.LEFT)
head_is.text_frame.margin_left = In(0.22)
tf = txbox(s, 0.82, 2.2, 5.5, 3.0)
for t in ["A guided symptom intake and triage assistant",
          "An explainer of likely conditions, in plain language",
          "A personal vault for reports, history and allergies",
          "A medication reminder and adherence tracker",
          "A router — it tells you how urgently to seek real care"]:
    bullet(tf, t, size=13, first=(t.startswith("A guided")), space_after=20)

head_not = box(s, 6.87, 1.55, 5.85, 0.46, fill=RGBColor(0xFF, 0xF1, 0xF2), edge=RGBColor(0xFE, 0xCD, 0xD3), radius=0.14)
label(head_not, "✕   IT IS NOT", size=13, color=ROSE, align=PP_ALIGN.LEFT)
head_not.text_frame.margin_left = In(0.22)
tf = txbox(s, 7.07, 2.2, 5.5, 3.0)
for t in ["A diagnosis — output is possibility, never verdict",
          "A prescription engine — no drug names, no dosages",
          "An emergency service — it escalates, it does not treat",
          "A replacement for a licensed clinician",
          "A store of record for legal or insurance purposes"]:
    bullet(tf, t, size=13, mark="▪", mark_color=ROSE, first=(t.startswith("A diagnosis")), space_after=20)

g = box(s, 0.62, 5.42, 12.1, 1.28, fill=AMBER_BG, edge=RGBColor(0xFE, 0xD7, 0xAA), radius=0.09)
g.text_frame.text = ""
tf = txbox(s, 0.92, 5.62, 11.5, 1.0)
para(tf, "⚠   THE GUARDRAIL PRINCIPLE", size=11, color=AMBER, bold=True, first=True, space_after=6)
para(tf, "Every AI response is wrapped with a confidence band, a mandatory 'consult a physician' notice, "
         "and an urgency level. Red-flag symptom combinations bypass the AI entirely and trigger immediate "
         "escalation — a deterministic rule, not a probabilistic guess.",
     size=12.5, color=SLATE, line=1.3, space_after=0)

# ================================================================ 04 MODULES
s = slide("Functional Modules", "Scope — Version 1.0",
          notes="MEMBER 1 closes here and hands to MEMBER 2.\nSix modules, all in scope for V1. Note the "
                "numbering — it maps one-to-one onto the Business Logic Layer components Member 2 shows next.")
mods = [
    ("01", "User Authentication\n& Profile", "Register / login, JWT sessions, hashed credentials, "
     "profile with age, sex, allergies, chronic conditions."),
    ("02", "Symptom Analysis", "Structured intake → red-flag rules → AI reasoning → ranked "
     "possible conditions with urgency level."),
    ("03", "AI Health Chatbot", "Follow-up conversation that remembers the session context and the "
     "user's stored history."),
    ("04", "Medical Report Upload", "PDF / image upload, OCR text extraction, AI summary of key values, "
     "attached to the timeline."),
    ("05", "Medicine Reminder", "Schedule doses, recurring jobs, push / e-mail notification, "
     "taken-or-missed adherence log."),
    ("06", "Medical History", "Unified chronological timeline of sessions, reports and medication, "
     "exportable for a real consultation."),
]
for i, (num, head, body) in enumerate(mods):
    x = 0.62 + (i % 3) * 4.1
    y = 1.62 + (i // 3) * 2.48
    c = box(s, x, y, 3.85, 2.2, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    chip = box(s, x + 0.28, y + 0.26, 0.52, 0.36, fill=MINT_SOFT, edge=None, radius=0.25)
    label(chip, num, size=11.5, color=TEAL_DARK)
    tf = txbox(s, x + 0.28, y + 0.74, 3.3, 1.3)
    para(tf, head, size=14, color=INK, bold=True, first=True, space_after=7, line=1.15)
    para(tf, body, size=11.5, color=SLATE, line=1.3, space_after=0)

# ================================================================ 05 TECH STACK
s = slide("Technology Stack — and Why", "Technologies",
          notes="MEMBER 2 begins.\n'Every row has a *why*. In an architecture review the choice matters less "
                "than the justification.' Be ready for: 'Why Node and not Java/Python?' — our workload is "
                "I/O-bound: most request time is spent waiting on the AI API and the database, not computing. "
                "Node's non-blocking event loop handles many concurrent waiting requests on modest hardware, "
                "and one JavaScript language across the stack suits a four-person team on a semester deadline.")
table(s, ["Layer", "Technology", "Why this choice"],
      [["Presentation", "React.js, Axios, Tailwind CSS",
        "Component reuse across 6 modules; SPA keeps chat responsive without full reloads"],
       ["API / Gateway", "Node.js + Express.js",
        "Non-blocking I/O suits an AI-call-heavy, I/O-bound workload; middleware chain for cross-cutting concerns"],
       ["Business Logic", "Express service modules",
        "Plain, testable JS services — triage rules, orchestration, scheduling; no framework lock-in"],
       ["AI Service", "Gemini / OpenAI behind an adapter",
        "Provider swapped by config, never by code change; enables fallback and cost comparison"],
       ["Database", "MongoDB (Atlas)",
        "Symptom sessions and reports are nested, evolving documents — a poor fit for rigid relational rows"],
       ["Cache / Jobs", "Redis + node-cron",
        "Caches repeated AI answers; drives reminder scheduling and OCR off the request path"],
       ["File Storage", "Cloud object storage (S3 / Cloudinary)",
        "Binary reports do not belong in a database; CDN delivery and signed URLs"],
       ["Security", "JWT, bcrypt, Helmet, express-validator",
        "Stateless auth for horizontal scaling; hashing, headers and validation at the edge"],
       ["Deployment", "Vercel (client) + Render / Cloud Run (API)",
        "Managed CI/CD from Git, HTTPS by default, environment-scoped secrets"]],
      l=0.62, t=1.5, w=12.1, col_w=[2.0, 3.5, 6.6], fs=11, hfs=11, row_h=0.53, head_h=0.4)

# ================================================================ 06 LAYERED DIAGRAM
s = slide("Architectural Style: Layered (N-Tier)", "Logical View",
          notes="MEMBER 2 — core slide.\n'We chose a layered architecture. Two rules make it real: strict "
                "downward dependency — a layer may only call the one beneath it — and no layer may skip "
                "a layer. The React client physically cannot reach MongoDB; it does not even know MongoDB "
                "exists. Cross-cutting concerns run vertically as Express middleware.'\n\nIf asked why not "
                "microservices: see our ADR-01 slide.")
layers = [
    ("PRESENTATION LAYER", "React SPA  ·  Components, Router, State, Axios client, Form validation",
     RGBColor(0xE0, 0xF2, 0xFE), RGBColor(0x07, 0x5A, 0x8C)),
    ("API / GATEWAY LAYER", "Express routes  ·  REST endpoints, JWT verify, rate limit, schema validation, error envelope",
     RGBColor(0xCC, 0xFB, 0xF1), TEAL_DARK),
    ("BUSINESS LOGIC LAYER", "Triage rules engine  ·  Session orchestrator  ·  Reminder scheduler  ·  History service",
     RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x15, 0x80, 0x3D)),
    ("AI SERVICE LAYER", "Provider adapter  ·  Prompt templates  ·  PII redaction  ·  Response validator  ·  Circuit breaker",
     RGBColor(0xED, 0xE9, 0xFE), INDIGO),
    ("DATA ACCESS LAYER", "Repositories  ·  Mongoose models  ·  Redis cache  ·  Object storage client",
     RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E)),
]
top = 1.55
for i, (name, desc, fill, col) in enumerate(layers):
    y = top + i * 0.96
    b = box(s, 0.62, y, 9.6, 0.80, fill=fill, edge=None, radius=0.1)
    b.text_frame.text = ""
    tf = txbox(s, 0.95, y + 0.13, 9.0, 0.58)
    para(tf, name, size=13, color=col, bold=True, first=True, space_after=3)
    para(tf, desc, size=10.5, color=SLATE, space_after=0)
    if i < len(layers) - 1:
        a = arrow(s, 5.28, y + 0.82, 0.28, 0.12, "down", color=RGBColor(0xCB, 0xD5, 0xE1))
        a.text_frame.text = ""

cc = box(s, 10.44, 1.55, 2.28, 4.64, fill=RGBColor(0xF8, 0xFA, 0xFC), edge=LINE, radius=0.06)
cc.text_frame.text = ""
tf = txbox(s, 10.64, 1.75, 1.9, 4.3)
para(tf, "CROSS-CUTTING\nCONCERNS", size=10, color=INK, bold=True, first=True, space_after=12, line=1.2)
for t in ["Authentication", "Authorization", "Input validation", "Logging & audit",
          "Error handling", "Configuration", "Caching", "Observability"]:
    bullet(tf, t, size=10.5, color=SLATE, mark="•", space_after=7)

note = box(s, 0.62, 6.36, 12.1, 0.44, fill=WHITE, edge=LINE, radius=0.1)
label(note, "Rule of the style:  each layer calls only the layer directly below it — no skipping, "
            "no upward calls.  →  swap MongoDB for MySQL and only the Data Layer changes.",
      size=11.5, color=SLATE, bold=False)

# ================================================================ 07 LAYER CONTRACTS
s = slide("Layer Responsibilities & Contracts", "Logical View",
          notes="MEMBER 2.\nRead the 'must never' column aloud — that is what enforces the architecture. "
                "A layered diagram anyone can draw; the discipline is in the constraints.")
table(s, ["Layer", "Owns", "Talks to", "Must never"],
      [["Presentation", "Rendering, client-side validation, session token storage, UX state",
        "API Layer over HTTPS/JSON only", "Hold business rules or database credentials"],
       ["API / Gateway", "Routing, authN/authZ, request schema validation, rate limiting, error shaping",
        "Business Logic Layer", "Contain medical logic or query the database directly"],
       ["Business Logic", "Triage rules, orchestration, scheduling, history assembly, policy",
        "AI Service + Data Access", "Know about HTTP, req/res objects or the AI vendor"],
       ["AI Service", "Prompt building, redaction, provider calls, retries, response validation",
        "External AI provider; returns DTOs", "Persist data or expose vendor-specific types upward"],
       ["Data Access", "CRUD, queries, indexing, caching, file persistence",
        "MongoDB, Redis, object storage", "Apply business rules or return raw driver objects"]],
      l=0.62, t=1.5, w=12.1, col_w=[1.7, 4.1, 2.6, 3.6], fs=11, hfs=11, row_h=0.86, head_h=0.4)

# ================================================================ 08 COMPONENT VIEW
s = slide("Component View", "Development View",
          notes="MEMBER 2 closes here.\n'This is the folder structure we will actually build. Each business "
                "component is a self-contained module — that is deliberate: it gives us clean seams if we "
                "ever extract a service later, which is exactly what ADR-01 depends on.'")
cols = [
    ("CLIENT  ·  React", RGBColor(0xE0, 0xF2, 0xFE), RGBColor(0x07, 0x5A, 0x8C),
     ["AuthPages  (login / register)", "SymptomIntakeForm", "ChatWindow", "ReportUploader",
      "ReminderManager", "HistoryTimeline", "ApiClient  (Axios + interceptors)",
      "AuthContext  (token, guard routes)"]),
    ("SERVER  ·  Express", MINT, TEAL_DARK,
     ["authController + authService", "symptomController", "TriageRulesEngine  ★",
      "AiOrchestrator", "reportService + OcrWorker", "reminderService + CronScheduler",
      "historyService", "middleware/  auth, validate, error, limit"]),
    ("DATA & EXTERNAL", RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E),
     ["UserRepository", "SymptomSessionRepository", "ReportRepository", "ReminderRepository",
      "AuditLogRepository", "RedisCacheClient", "ObjectStorageClient", "AiProviderAdapter  → Gemini / OpenAI"]),
]
for i, (head, fill, col, items) in enumerate(cols):
    x = 0.62 + i * 4.1
    c = box(s, x, 1.52, 3.85, 4.5, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    hb = box(s, x, 1.52, 3.85, 0.5, fill=fill, edge=None, radius=0.1)
    label(hb, head, size=11.5, color=col)
    tf = txbox(s, x + 0.26, 2.2, 3.35, 3.7)
    for j, it in enumerate(items):
        star = "★" in it
        bullet(tf, it.replace("  ★", ""), size=11,
               color=(TEAL_DARK if star else SLATE), bold=star,
               mark="▪", mark_color=(ROSE if star else MUTED),
               first=(j == 0), space_after=10.5)

n = box(s, 0.62, 6.24, 12.1, 0.44, fill=MINT_SOFT, edge=MINT, radius=0.1)
label(n, "★  TriageRulesEngine is deterministic and runs before any AI call — the safety core of the system.",
      size=11, color=TEAL_DARK, bold=False)

# ================================================================ 09 WORKFLOW
s = slide("Runtime Workflow — Symptom Analysis", "Process View",
          notes="MEMBER 3 begins.\n'This is the end-to-end path of a single request.' Walk 1→8 in order. "
                "The two callouts matter: step 3 can terminate the whole flow without ever calling the AI, "
                "and step 7 writes to the database BEFORE responding, so history is never lost if the client "
                "drops. Target: p95 under 3 seconds.")
steps = [
    ("1", "User", "Logs in, enters\nsymptoms or\nuploads a report", RGBColor(0xE0, 0xF2, 0xFE), RGBColor(0x07, 0x5A, 0x8C)),
    ("2", "React Client", "Client-side validation,\nattaches JWT,\nPOST /api/v1/...", MINT, TEAL_DARK),
    ("3", "Express API", "Verify token,\nrate limit,\nschema validation", MINT, TEAL_DARK),
    ("4", "Triage Engine", "Deterministic\nred-flag rules\n→ may escalate now", RGBColor(0xFF, 0xE4, 0xE6), ROSE),
    ("5", "AI Service", "Redact PII, build\nprompt, call provider,\nvalidate response", RGBColor(0xED, 0xE9, 0xFE), INDIGO),
    ("6", "Database", "Persist session,\nupdate history,\nwrite audit log", RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E)),
    ("7", "Response", "Conditions, urgency,\nconfidence, advice,\ndisclaimer", RGBColor(0xDC, 0xFC, 0xE7), RGBColor(0x15, 0x80, 0x3D)),
]
x0, bw, gap = 0.62, 1.52, 0.22
for i, (num, head, body, fill, col) in enumerate(steps):
    x = x0 + i * (bw + gap)
    c = box(s, x, 1.95, bw, 1.85, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    hb = box(s, x, 1.95, bw, 0.44, fill=fill, edge=None, radius=0.14)
    label(hb, head, size=10.5, color=col)
    chip = box(s, x + bw / 2 - 0.18, 1.52, 0.36, 0.36, fill=col, edge=None, radius=0.5)
    label(chip, num, size=11, color=WHITE)
    tf = txbox(s, x + 0.12, 2.54, bw - 0.24, 1.2)
    para(tf, body, size=9.5, color=SLATE, first=True, space_after=0, line=1.25, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        a = arrow(s, x + bw + 0.03, 2.79, 0.16, 0.18, "right", color=RGBColor(0xCB, 0xD5, 0xE1))
        a.text_frame.text = ""

c1 = box(s, 0.62, 4.2, 5.85, 1.15, fill=RGBColor(0xFF, 0xF1, 0xF2), edge=RGBColor(0xFE, 0xCD, 0xD3), radius=0.1)
c1.text_frame.text = ""
tf = txbox(s, 0.9, 4.4, 5.3, 0.9)
para(tf, "SHORT-CIRCUIT PATH", size=10, color=ROSE, bold=True, first=True, space_after=5)
para(tf, "If step 4 matches a red flag (e.g. chest pain + breathlessness), the AI is never called — "
         "the user is told to seek emergency care immediately.", size=11.5, color=SLATE, line=1.28, space_after=0)

c2 = box(s, 6.87, 4.2, 5.85, 1.15, fill=MINT_SOFT, edge=MINT, radius=0.1)
c2.text_frame.text = ""
tf = txbox(s, 7.15, 4.4, 5.3, 0.9)
para(tf, "FAILURE PATH", size=10, color=TEAL_DARK, bold=True, first=True, space_after=5)
para(tf, "If the AI provider times out, the circuit breaker opens: fall back to the second provider, "
         "then to rules-only guidance. The system degrades, it does not fail.",
     size=11.5, color=SLATE, line=1.28, space_after=0)

b = box(s, 0.62, 5.65, 12.1, 0.6, fill=WHITE, edge=LINE, radius=0.1)
label(b, "Long-running work — OCR of uploaded reports, reminder dispatch — is queued and processed "
         "asynchronously, never inside the request/response cycle.", size=11.5, color=SLATE, bold=False)

# ================================================================ 10 AI PIPELINE
s = slide("Inside the AI Service Layer", "The Safety Pipeline",
          notes="MEMBER 3 — the strongest technical slide, spend time here.\n'A language model is "
                "probabilistic. Our answer is to surround it with deterministic code on both sides: "
                "everything before step 4 constrains what it sees, everything after constrains what it "
                "is allowed to say.'\n\nExpected question: 'What if it hallucinates a drug name?' → "
                "Step 6 rejects the response on a banned-content check and we return the rules-based "
                "guidance instead.")
pipe = [
    ("1", "Normalise", "Structured intake: symptom, duration,\nseverity, onset — not free text alone"),
    ("2", "Redact PII", "Name, phone, e-mail, ID stripped\nbefore anything leaves our boundary"),
    ("3", "Red-flag rules", "Deterministic table lookup.\nMatch → escalate, skip AI entirely"),
    ("4", "Assemble context", "Age, sex, allergies, chronic conditions,\nrecent history + prompt template"),
    ("5", "Provider call", "Adapter → Gemini / OpenAI.\nTimeout, retry, circuit breaker, fallback"),
    ("6", "Validate output", "Schema check + banned-content check:\nno dosages, no prescriptions"),
    ("7", "Wrap & persist", "Attach confidence, urgency, disclaimer.\nSave to history, write audit log"),
]
for i, (n_, head, body) in enumerate(pipe):
    y = 1.5 + i * 0.755
    chip = box(s, 0.62, y + 0.09, 0.42, 0.42, fill=TEAL_DARK, edge=None, radius=0.5)
    label(chip, n_, size=11, color=WHITE)
    hb = box(s, 1.22, y, 2.55, 0.6, fill=(RGBColor(0xFF, 0xE4, 0xE6) if i == 2 else MINT_SOFT),
             edge=(RGBColor(0xFE, 0xCD, 0xD3) if i == 2 else MINT), radius=0.12)
    label(hb, head, size=12, color=(ROSE if i == 2 else TEAL_DARK), align=PP_ALIGN.LEFT)
    hb.text_frame.margin_left = In(0.16)
    tf = txbox(s, 4.0, y + 0.06, 4.5, 0.55)
    para(tf, body, size=10.5, color=SLATE, first=True, space_after=0, line=1.22)
    if i < len(pipe) - 1:
        a = arrow(s, 0.75, y + 0.56, 0.16, 0.16, "down", color=RGBColor(0xCB, 0xD5, 0xE1))
        a.text_frame.text = ""

p = box(s, 8.85, 1.5, 3.87, 4.95, fill=RGBColor(0xF8, 0xFA, 0xFC), edge=LINE)
p.text_frame.text = ""
tf = txbox(s, 9.1, 1.75, 3.4, 4.5)
para(tf, "WHY THIS SHAPE", size=10, color=INK, bold=True, first=True, space_after=12)
for h, b_ in [("Deterministic before probabilistic",
               "Safety-critical decisions never depend on a model's mood."),
              ("Vendor-neutral by construction",
               "The adapter is the only file that knows Gemini from OpenAI."),
              ("Fail-soft, never fail-open",
               "Provider down → degraded rules-only advice, still safe."),
              ("Privacy at the boundary",
               "Identifiers never reach a third-party API."),
              ("Auditable",
               "Every AI input and output is logged for review.")]:
    para(tf, h, size=11.5, color=TEAL_DARK, bold=True, space_after=3)
    para(tf, b_, size=10.5, color=SLATE, space_after=13, line=1.25)

# ================================================================ 11 DATA MODEL
s = slide("Data Model", "Data View  ·  MongoDB Collections",
          notes="MEMBER 3.\n'Six collections. Note we store a *denormalised snapshot* of the AI response — "
                "if we change prompts or providers next month, old sessions must still render exactly as the "
                "user saw them. That is an architectural choice about immutability of history, not laziness.'\n\n"
                "If asked why not MySQL → ADR-03.")
colls = [
    ("users", TEAL_DARK, MINT_SOFT, ["_id, email (unique idx)", "passwordHash  (bcrypt)", "profile: age, sex, allergies[]",
                                     "chronicConditions[]", "createdAt, lastLoginAt"]),
    ("symptom_sessions", INDIGO, RGBColor(0xED, 0xE9, 0xFE), ["_id, userId  (idx)", "input: symptoms[], duration, severity",
                                                              "triageResult: urgency, redFlags[]", "aiResponse  (immutable snapshot)",
                                                              "confidence, createdAt"]),
    ("reports", RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFE, 0xF3, 0xC7), ["_id, userId  (idx)", "fileUrl, mimeType, sizeBytes",
                                                                         "extractedText  (OCR)", "aiSummary, keyValues{}", "uploadedAt"]),
    ("reminders", RGBColor(0x15, 0x80, 0x3D), RGBColor(0xDC, 0xFC, 0xE7), ["_id, userId  (idx)", "medicineName, dosageNote",
                                                                           "schedule: cron / times[]", "channel, active",
                                                                           "adherenceLog[]  {at, taken}"]),
    ("chat_messages", RGBColor(0x07, 0x5A, 0x8C), RGBColor(0xE0, 0xF2, 0xFE), ["_id, userId, sessionId  (idx)", "role: user | assistant",
                                                                               "content, tokensUsed", "createdAt  (TTL optional)"]),
    ("audit_logs", ROSE, RGBColor(0xFF, 0xE4, 0xE6), ["_id, userId, action", "resource, resourceId", "ipHash, userAgent",
                                                      "aiProvider, latencyMs", "timestamp"]),
]
for i, (name, col, fill, fields) in enumerate(colls):
    x = 0.62 + (i % 3) * 4.1
    y = 1.5 + (i // 3) * 2.52
    c = box(s, x, y, 3.85, 2.24, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    hb = box(s, x, y, 3.85, 0.46, fill=fill, edge=None, radius=0.12)
    label(hb, name, size=12, color=col, align=PP_ALIGN.LEFT)
    hb.text_frame.margin_left = In(0.22)
    tf = txbox(s, x + 0.26, y + 0.62, 3.35, 1.55)
    for j, f in enumerate(fields):
        bullet(tf, f, size=10, color=SLATE, mark="•", mark_color=MUTED,
               first=(j == 0), space_after=7, mar=0.15)

n = box(s, 0.62, 6.36, 12.1, 0.44, fill=MINT_SOFT, edge=MINT, radius=0.1)
label(n, "Relationships by userId reference  ·  indexed on userId + createdAt for timeline queries  ·  "
         "audit_logs is append-only", size=11, color=TEAL_DARK, bold=False)

# ================================================================ 12 API DESIGN
s = slide("REST API Design", "Interface View",
          notes="MEMBER 3 closes here and hands to MEMBER 4.\n'Resource-oriented URLs, HTTP verbs carry the "
                "action, versioned prefix so we can evolve without breaking the client. Every error uses one "
                "envelope so the React layer has exactly one error path to handle.'")
table(s, ["Method", "Endpoint", "Purpose", "Auth"],
      [["POST", "/api/v1/auth/register", "Create account, hash password, issue JWT", "Public"],
       ["POST", "/api/v1/auth/login", "Authenticate, return access + refresh token", "Public"],
       ["POST", "/api/v1/symptoms/analyze", "Run triage → AI pipeline, persist session", "JWT"],
       ["GET", "/api/v1/history?page=&limit=", "Paginated unified timeline for the user", "JWT"],
       ["POST", "/api/v1/reports", "Multipart upload, queue OCR job", "JWT"],
       ["GET", "/api/v1/reports/:id", "Fetch report metadata, summary, signed file URL", "JWT + owner"],
       ["POST", "/api/v1/chat/messages", "Send follow-up message in a session context", "JWT"],
       ["GET / POST / PUT / DELETE", "/api/v1/reminders/:id?", "Full CRUD over medication reminders", "JWT + owner"]],
      l=0.62, t=1.5, w=12.1, col_w=[2.3, 3.6, 5.2, 1.6], fs=11, hfs=11, row_h=0.47, head_h=0.4)

conv = box(s, 0.62, 5.82, 12.1, 1.0, fill=RGBColor(0xF8, 0xFA, 0xFC), edge=LINE, radius=0.1)
conv.text_frame.text = ""
tf = txbox(s, 0.92, 5.98, 11.5, 0.8)
para(tf, "API CONVENTIONS", size=10, color=INK, bold=True, first=True, space_after=7)
rich(tf, [("Versioned prefix ", INK, True), ("/api/v1  ·  ", SLATE, False),
          ("JSON only  ·  ", SLATE, False), ("Bearer JWT ", INK, True),
          ("in Authorization header  ·  ", SLATE, False),
          ("Uniform error envelope ", INK, True),
          ("{ success, data, error: { code, message } }  ·  ", SLATE, False),
          ("Correct status codes ", INK, True),
          ("(200 / 201 / 400 / 401 / 403 / 404 / 429 / 503)  ·  ", SLATE, False),
          ("Pagination ", INK, True), ("on every list endpoint", SLATE, False)], size=11.5, line=1.35)

# ================================================================ 13 ADRs
s = slide("Key Architectural Decisions", "ADRs  ·  Decision Log",
          notes="MEMBER 4 begins — this is the highest-value slide for marks.\n'For each decision we "
                "recorded the alternative we rejected and the price we accept.'\n\nRehearse ADR-01: "
                "microservices would give independent scaling and deployment, but four students on a "
                "semester deadline cannot absorb distributed tracing, network failure handling and "
                "multi-service DevOps. We chose a MODULAR monolith — clean module seams so the AI "
                "service can be extracted later without a rewrite. That is deferring a decision, not avoiding it.")
table(s, ["ID", "Decision", "Rationale", "Trade-off accepted"],
      [["ADR-01", "Modular monolith, not microservices",
        "One deployable unit fits a 4-person team and a semester timeline; module seams preserved",
        "Whole app scales together; extraction work deferred"],
       ["ADR-02", "Deterministic triage engine in front of the AI",
        "Safety-critical escalation must not depend on a probabilistic model",
        "Rules table must be curated and maintained manually"],
       ["ADR-03", "MongoDB over MySQL",
        "Sessions and reports are nested, schema-evolving documents; no heavy joins in our workload",
        "Weaker cross-collection integrity; enforced in application code"],
       ["ADR-04", "AI provider behind an Adapter interface",
        "Avoids vendor lock-in; enables fallback, A/B cost comparison and offline mocking in tests",
        "One extra indirection layer to build and maintain"],
       ["ADR-05", "Stateless JWT authentication",
        "No server session store → any instance serves any request; horizontal scaling is trivial",
        "Token revocation needs a short TTL plus refresh-token rotation"],
       ["ADR-06", "Async queue for OCR and reminders",
        "Keeps p95 request latency low; slow work never blocks the user's thread",
        "Adds a worker process and eventual-consistency semantics"]],
      l=0.62, t=1.5, w=12.1, col_w=[0.95, 3.0, 4.9, 3.25], fs=10.5, hfs=11, row_h=0.79, head_h=0.4)

# ================================================================ 14 QUALITY ATTRIBUTES
s = slide("Quality Attributes & Architectural Tactics", "Non-Functional Requirements",
          notes="MEMBER 4.\n'Functional requirements say what the system does; quality attributes decide "
                "what the architecture looks like. Each row is a measurable scenario paired with the tactic "
                "that achieves it — not an adjective.'\n\nIf asked which dominates: safety first, then "
                "privacy. We would trade latency for either without hesitating.")
table(s, ["Attribute", "Measurable scenario", "Architectural tactic"],
      [["Safety", "100% of defined red-flag combinations escalate, even if the AI is unavailable",
        "Deterministic rules engine evaluated before the AI call; rules-only fallback path"],
       ["Privacy", "No direct identifier ever leaves our trust boundary",
        "PII redaction stage in the AI layer; encryption at rest; append-only audit log"],
       ["Performance", "p95 response for symptom analysis under 3 s at 100 concurrent users",
        "Redis cache for repeated queries; async offload of OCR; indexed queries; streaming responses"],
       ["Availability", "Core triage stays usable when a single AI provider fails",
        "Circuit breaker, provider fallback, graceful degradation to rules-only mode"],
       ["Modifiability", "Swap AI vendor or database with changes confined to one layer",
        "Adapter and Repository patterns; strict layer dependency rule; DTOs at every boundary"],
       ["Scalability", "Handle traffic growth without re-architecting",
        "Stateless API instances behind a load balancer; managed DB with connection pooling"],
       ["Security", "Withstand credential stuffing and injection attempts",
        "bcrypt hashing, rate limiting, schema validation at the edge, Helmet headers, RBAC"],
       ["Usability", "A non-technical user completes an analysis in under 90 seconds",
        "Guided structured intake instead of free-text; plain-language output with urgency colour coding"]],
      l=0.62, t=1.5, w=12.1, col_w=[1.6, 4.6, 5.9], fs=10.5, hfs=11, row_h=0.6, head_h=0.4)

# ================================================================ 15 PATTERNS
s = slide("Design Patterns Applied", "Architecture & Design",
          notes="MEMBER 4.\nDo not just name the patterns — name where each one lives in OUR code. "
                "Examiners test that. Strongest three to emphasise: Adapter (AI vendor), Repository "
                "(database swap), Circuit Breaker (resilience).")
pats = [
    ("Layered", "Overall system decomposition into five stacked layers", TEAL_DARK, MINT_SOFT),
    ("Adapter", "AiProviderAdapter — one interface, Gemini or OpenAI behind it", INDIGO, RGBColor(0xED, 0xE9, 0xFE)),
    ("Repository", "UserRepository, ReportRepository — hides Mongoose from services", RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFE, 0xF3, 0xC7)),
    ("Strategy", "Interchangeable prompt strategies per symptom category", RGBColor(0x15, 0x80, 0x3D), RGBColor(0xDC, 0xFC, 0xE7)),
    ("Chain of Responsibility", "Express middleware pipeline: auth → validate → limit → handle", RGBColor(0x07, 0x5A, 0x8C), RGBColor(0xE0, 0xF2, 0xFE)),
    ("Facade", "AiOrchestrator exposes one call over a seven-stage pipeline", TEAL_DARK, MINT_SOFT),
    ("Observer / Pub-Sub", "Reminder events fan out to notification channels", ROSE, RGBColor(0xFF, 0xE4, 0xE6)),
    ("Circuit Breaker", "Trips after repeated AI provider failures; opens the fallback path", INDIGO, RGBColor(0xED, 0xE9, 0xFE)),
    ("DTO", "Typed objects crossing every layer boundary — no leaking internals", RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFE, 0xF3, 0xC7)),
]
for i, (name, desc, col, fill) in enumerate(pats):
    x = 0.62 + (i % 3) * 4.1
    y = 1.55 + (i // 3) * 1.68
    c = box(s, x, y, 3.85, 1.42, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    dot = box(s, x + 0.26, y + 0.26, 0.28, 0.28, fill=fill, edge=None, radius=0.5)
    dot.text_frame.text = ""
    d2 = box(s, x + 0.335, y + 0.335, 0.13, 0.13, fill=col, edge=None, radius=0.5)
    d2.text_frame.text = ""
    tf = txbox(s, x + 0.66, y + 0.24, 2.95, 1.0)
    para(tf, name, size=12.5, color=INK, bold=True, first=True, space_after=5)
    para(tf, desc, size=10.5, color=SLATE, line=1.25, space_after=0)

n = box(s, 0.62, 6.5, 12.1, 0.4, fill=MINT_SOFT, edge=MINT, radius=0.1)
label(n, "Each pattern is justified by a quality attribute — Adapter serves modifiability, "
         "Circuit Breaker serves availability, Repository serves testability.", size=11, color=TEAL_DARK, bold=False)

# ================================================================ 16 SECURITY
s = slide("Security & Privacy Architecture", "Cross-Cutting Concern",
          notes="MEMBER 4.\n'Health data is among the most sensitive category of personal data. We treat "
                "the third-party AI provider as an UNTRUSTED external party — that is why redaction sits "
                "at our boundary and not at theirs.'\n\nIf asked about compliance: we are not claiming "
                "HIPAA or DPDP certification for a student project, but the design follows their core "
                "principles — minimisation, encryption, access control, auditability.")
groups = [
    ("IDENTITY", TEAL_DARK, MINT_SOFT,
     ["bcrypt password hashing with salt rounds", "JWT access token, short TTL", "Refresh-token rotation on use",
      "Role-based access: user / admin", "Ownership check on every resource"]),
    ("TRANSPORT & EDGE", RGBColor(0x07, 0x5A, 0x8C), RGBColor(0xE0, 0xF2, 0xFE),
     ["HTTPS / TLS enforced end to end", "Helmet security headers", "Strict CORS allow-list",
      "Rate limiting per IP and per user", "Schema validation before any handler"]),
    ("DATA", RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFE, 0xF3, 0xC7),
     ["Encryption at rest on the managed DB", "Signed, expiring URLs for report files",
      "Secrets in environment variables only", "Append-only audit log of all access", "Right to delete: full account purge"]),
    ("AI BOUNDARY", INDIGO, RGBColor(0xED, 0xE9, 0xFE),
     ["PII redacted before every provider call", "No raw report images sent to the AI",
      "Prompt-injection filtering on user text", "Output validated against a schema",
      "Provider traffic logged for audit"]),
]
for i, (head, col, fill, items) in enumerate(groups):
    x = 0.62 + i * 3.08
    c = box(s, x, 1.52, 2.83, 4.4, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    hb = box(s, x, 1.52, 2.83, 0.46, fill=fill, edge=None, radius=0.12)
    label(hb, head, size=10.5, color=col)
    tf = txbox(s, x + 0.22, 2.16, 2.42, 3.6)
    for j, it in enumerate(items):
        bullet(tf, it, size=10.5, color=SLATE, mark="▪", mark_color=col,
               first=(j == 0), space_after=13, mar=0.17)

n = box(s, 0.62, 6.14, 12.1, 0.5, fill=AMBER_BG, edge=RGBColor(0xFE, 0xD7, 0xAA), radius=0.1)
label(n, "⚠   Design principle: the AI provider is treated as an untrusted external party. "
         "Identifiable data never crosses that boundary.", size=11.5, color=AMBER, bold=False)

# ================================================================ 17 DEPLOYMENT
s = slide("Deployment View", "Physical View",
          notes="MEMBER 4 (or MEMBER 2 if you rebalance).\n'Managed services throughout — as a student "
                "team our scarcest resource is time, so we buy infrastructure rather than operate it. "
                "The API tier is stateless, which is what lets us add instances behind the load balancer "
                "with no session affinity.'")
GAP = 0.267
tiers = [
    (2.95, "CLIENT TIER", RGBColor(0xE0, 0xF2, 0xFE), RGBColor(0x07, 0x5A, 0x8C),
     ["Browser / mobile web", "React SPA on Vercel", "Global CDN + HTTPS"], "HTTPS / REST"),
    (2.95, "APPLICATION TIER", MINT, TEAL_DARK,
     ["Node + Express (Render)", "Stateless — N instances", "Load balancer, auto-scale",
      "Worker: OCR + cron jobs"], "TLS / driver"),
    (2.95, "DATA TIER", RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0x92, 0x40, 0x0E),
     ["MongoDB Atlas (managed)", "Redis — cache + queue", "Object storage — reports"], "HTTPS / SDK"),
    (2.45, "EXTERNAL", RGBColor(0xED, 0xE9, 0xFE), INDIGO,
     ["Gemini API (primary)", "OpenAI API (fallback)", "E-mail / push provider"], None),
]
x = 0.62
for i, (w_, head, fill, col, items, link) in enumerate(tiers):
    c = box(s, x, 1.75, w_, 3.2, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    hb = box(s, x, 1.75, w_, 0.48, fill=fill, edge=None, radius=0.12)
    label(hb, head, size=11, color=col)
    tf = txbox(s, x + 0.22, 2.42, w_ - 0.44, 2.4)
    for j, it in enumerate(items):
        bullet(tf, it, size=10.5, color=SLATE, mark="▪", mark_color=col,
               first=(j == 0), space_after=14, mar=0.17)
    if link:
        a = arrow(s, x + w_ + 0.045, 3.25, 0.18, 0.2, "right", color=RGBColor(0xCB, 0xD5, 0xE1))
        a.text_frame.text = ""
        tf = txbox(s, x + w_ + GAP / 2 - 0.7, 5.04, 1.4, 0.28)
        para(tf, link, size=8.5, color=MUTED, first=True, space_after=0, align=PP_ALIGN.CENTER)
    x += w_ + GAP

foot = box(s, 0.62, 5.5, 12.1, 1.2, fill=RGBColor(0xF8, 0xFA, 0xFC), edge=LINE, radius=0.1)
foot.text_frame.text = ""
tf = txbox(s, 0.92, 5.7, 11.5, 1.0)
para(tf, "DELIVERY PIPELINE", size=10, color=INK, bold=True, first=True, space_after=7)
rich(tf, [("Git push  →  ", SLATE, False), ("automated build & tests  →  ", INK, True),
          ("preview deployment per branch  →  ", SLATE, False), ("review  →  ", INK, True),
          ("promote to production.  ", SLATE, False),
          ("Environment-scoped secrets; no credentials in the repository; rollback is a one-click "
           "redeploy of the previous build.", SLATE, False)], size=11.5, line=1.35)

# ================================================================ 18 RISKS
s = slide("Risks & Mitigations", "Architectural Risk Register",
          notes="MEMBER 4.\n'Naming your risks before the panel does is the point of this slide.' The "
                "honest one is row 1 — we cannot eliminate hallucination, so we CONTAIN it: constrain "
                "the input, validate the output, and never let the model own the escalation decision.")
table(s, ["Risk", "Impact", "Mitigation"],
      [["AI produces an incorrect or hallucinated suggestion", "High · Safety",
        "Rules engine owns escalation; output schema + banned-content validation; confidence shown; mandatory disclaimer"],
       ["User treats output as a diagnosis and skips real care", "High · Safety",
        "Urgency-first UI, persistent disclaimer, explicit 'see a doctor' call-to-action on every result"],
       ["AI provider outage, rate limit or cost spike", "Medium · Availability",
        "Circuit breaker, secondary provider fallback, Redis response cache, per-user quotas"],
       ["Breach or leakage of sensitive health data", "High · Privacy",
        "PII redaction at the boundary, encryption at rest and in transit, RBAC, audit logging, least privilege"],
       ["Latency exceeds user patience on AI calls", "Medium · Performance",
        "Streamed partial responses, cached common queries, async offload of OCR, p95 budget tracked"],
       ["Prompt injection through free-text symptom input", "Medium · Security",
        "Structured intake over free text, input sanitisation, system-prompt isolation, output validation"],
       ["Scope creep across six modules in one semester", "Medium · Delivery",
        "Phased roadmap; modules 1–3 form the graded core, 4–6 follow incrementally"]],
      l=0.62, t=1.5, w=12.1, col_w=[3.5, 1.9, 6.7], fs=10.5, hfs=11, row_h=0.66, head_h=0.4)

# ================================================================ 19 ROADMAP
s = slide("Roadmap & Future Enhancements", "Evolution",
          notes="MEMBER 4.\n'Phase 1 is what we commit to delivering this semester. Phases 2 and 3 show the "
                "architecture has room to grow — and crucially, none of them require a rewrite, because "
                "of the module seams in ADR-01 and the adapter in ADR-04.'")
phases = [
    ("PHASE 1", "This semester — committed", TEAL_DARK, MINT_SOFT,
     ["Authentication & profile", "Symptom analysis + triage engine", "AI chatbot with session context",
      "Medical history timeline", "Deployed, documented, tested"]),
    ("PHASE 2", "Next iteration", INDIGO, RGBColor(0xED, 0xE9, 0xFE),
     ["Report upload + OCR extraction", "Medicine reminders & adherence", "Multilingual support (regional languages)",
      "Doctor appointment booking", "Nearby hospital & pharmacy lookup"]),
    ("PHASE 3", "Long term vision", RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFE, 0xF3, 0xC7),
     ["Emergency SOS with location sharing", "Wearable device integration (vitals)", "Extract AI layer into its own service",
      "Event-driven core for real-time alerts", "FHIR-compliant record exchange"]),
]
for i, (ph, sub, col, fill, items) in enumerate(phases):
    x = 0.62 + i * 4.1
    c = box(s, x, 1.6, 3.85, 4.3, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    hb = box(s, x, 1.6, 3.85, 0.72, fill=fill, edge=None, radius=0.12)
    label(hb, ph, size=12.5, color=col, align=PP_ALIGN.LEFT, sub=sub, sub_size=9.5, sub_color=SLATE)
    hb.text_frame.margin_left = In(0.24)
    tf = txbox(s, x + 0.26, 2.52, 3.35, 3.2)
    for j, it in enumerate(items):
        bullet(tf, it, size=11.5, color=SLATE, mark="▪", mark_color=col,
               first=(j == 0), space_after=15)

n = box(s, 0.62, 6.14, 12.1, 0.44, fill=MINT_SOFT, edge=MINT, radius=0.1)
label(n, "None of these require re-architecting — the module seams (ADR-01) and the provider adapter "
         "(ADR-04) were chosen precisely to absorb this growth.", size=11, color=TEAL_DARK, bold=False)

# ================================================================ 20 TEAM & CLOSE
s = slide("Team Responsibilities", "Work Split & Presentation Roles",
          notes="Whoever is closing.\nRead the split, then deliver the closing line and invite questions.\n\n"
                "Likely panel questions to rehearse:\n"
                "1. Why layered and not microservices? → ADR-01\n"
                "2. What if the AI is wrong? → rules engine + output validation + disclaimer\n"
                "3. Why MongoDB? → ADR-03, document shape and schema evolution\n"
                "4. How do you protect health data? → redaction at the boundary, encryption, audit log\n"
                "5. How does it scale? → stateless JWT + horizontal instances + cache + async queue")
members = [
    ("MEMBER 1", "Project Overview & Requirements", "Problem framing, scope and safety boundary, module "
     "definition, requirement documentation.", "Slides 1–4", TEAL_DARK, MINT_SOFT),
    ("MEMBER 2", "Architecture & Technology", "Layered design, layer contracts, component view, stack "
     "justification, deployment view.", "Slides 5–8, 17", INDIGO, RGBColor(0xED, 0xE9, 0xFE)),
    ("MEMBER 3", "Workflow, AI Pipeline & Data", "Runtime flow, AI safety pipeline, data model, "
     "REST API contract.", "Slides 9–12", RGBColor(0x15, 0x80, 0x3D), RGBColor(0xDC, 0xFC, 0xE7)),
    ("MEMBER 4", "Decisions, Quality & Risk", "ADR log, quality attributes and tactics, design patterns, "
     "security, risks, roadmap.", "Slides 13–16, 18–19", RGBColor(0x92, 0x40, 0x0E), RGBColor(0xFE, 0xF3, 0xC7)),
]
for i, (m, role, desc, slides_, col, fill) in enumerate(members):
    y = 1.5 + i * 1.14
    c = box(s, 0.62, y, 12.1, 1.0, fill=WHITE, edge=LINE)
    c.text_frame.text = ""
    bar = box(s, 0.62, y, 0.05, 1.0, fill=col, edge=None, radius=0.4)
    bar.text_frame.text = ""
    chip = box(s, 0.92, y + 0.3, 1.35, 0.4, fill=fill, edge=None, radius=0.2)
    label(chip, m, size=10, color=col)
    tf = txbox(s, 2.5, y + 0.2, 7.6, 0.7)
    para(tf, role, size=13, color=INK, bold=True, first=True, space_after=4)
    para(tf, desc, size=11, color=SLATE, space_after=0, line=1.2)
    tf = txbox(s, 10.4, y + 0.38, 2.1, 0.3)
    para(tf, slides_, size=10.5, color=MUTED, first=True, space_after=0, align=PP_ALIGN.RIGHT)

close = box(s, 0.62, 6.14, 12.1, 0.66, fill=TEAL_DEEP, edge=None, radius=0.1)
label(close, "MedIntel makes healthcare guidance faster and better organised — while keeping "
             "the clinician, not the algorithm, in charge of the diagnosis.",
      size=12.5, color=MINT, bold=False)

# ================================================================ 21 THANK YOU
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = TEAL_DEEP
for x, y, d, col in [(10.0, -1.2, 4.8, RGBColor(0x11, 0x5E, 0x59)), (-1.0, 4.8, 3.6, RGBColor(0x0E, 0x57, 0x53))]:
    c = box(s, x, y, d, d, fill=col, edge=None, shape=MSO_SHAPE.OVAL)
    c.text_frame.text = ""
tf = txbox(s, 1.0, 2.75, 11.0, 1.2)
para(tf, "Thank You", size=52, color=WHITE, bold=True, first=True, space_after=10, line=1.0)
para(tf, "Questions and feedback welcome.", size=17, color=MINT, space_after=0)
rule = box(s, 1.0, 4.5, 1.5, 0.035, fill=TEAL, edge=None, shape=MSO_SHAPE.RECTANGLE)
rule.text_frame.text = ""
tf = txbox(s, 1.0, 4.85, 11.0, 0.4)
para(tf, "MedIntel  ·  Software Architecture, Semester 5", size=12,
     color=RGBColor(0x7F, 0xD1, 0xC6), first=True, space_after=0)
s.notes_slide.notes_text_frame.text = (
    "Anticipated questions:\n"
    "• Why layered, not microservices? → ADR-01: team size + timeline; modular seams preserve the option.\n"
    "• What if the AI hallucinates? → Rules engine owns escalation; output schema + banned-content check.\n"
    "• Why MongoDB over MySQL? → ADR-03: nested, schema-evolving documents; no join-heavy queries.\n"
    "• Is this legal / regulated? → Positioned as decision SUPPORT with mandatory disclaimers; "
    "no prescriptions, no dosages.\n"
    "• How does it scale? → Stateless JWT, horizontal API instances, Redis cache, async job queue.\n"
    "• Biggest architectural risk? → Dependence on a third-party AI provider; mitigated by the "
    "adapter, circuit breaker and fallback provider.")

out = "MedIntel_Project_Overview.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
